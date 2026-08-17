from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.extract.base import AbstractExtractor, ExtractedDoc, Page


class PptxExtractor(AbstractExtractor):
    method = "python-pptx"

    def extract(self, path: Path) -> ExtractedDoc:
        prs = Presentation(str(path))
        pages: list[Page] = []
        deck_title: str | None = None
        for idx, slide in enumerate(prs.slides, start=1):
            title = ""
            body_chunks: list[str] = []
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = (slide.shapes.title.text or "").strip()
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = (shape.text_frame.text or "").strip()
                if text:
                    body_chunks.append(text)
            parts: list[str] = []
            if title:
                parts.append(f"# {title}")
                if deck_title is None:
                    deck_title = title
            parts.extend(body_chunks)

            notes_slide = slide.notes_slide if slide.has_notes_slide else None
            if notes_slide and notes_slide.notes_text_frame:
                note_text = (notes_slide.notes_text_frame.text or "").strip()
                if note_text:
                    quoted = "\n".join("> " + ln for ln in note_text.splitlines())
                    parts.append("**Speaker notes:**\n" + quoted)

            text_md = "\n\n".join(parts).strip()
            pages.append(Page(no=idx, text_md=text_md))

        word_count = sum(len(p.text_md.split()) for p in pages)
        return ExtractedDoc(
            pages=pages,
            title=deck_title,
            word_count=word_count,
            extraction_method=self.method,
        )
